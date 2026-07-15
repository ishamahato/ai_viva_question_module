from fastapi import FastAPI, Depends, HTTPException, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from sqlalchemy.orm import Session
from passlib.context import CryptContext
from jose import JWTError, jwt
from datetime import datetime, timedelta, timezone
from typing import List, Optional
import io, re

from database import Base, engine, get_db, settings
from database import Teacher, QuestionSet, Question
from database import (
    TeacherRegister, TeacherLogin, TeacherOut, Token,
    QuestionSetOut, QuestionSetDetail, QuestionSetDetailWithAnswers
)

# ─── App Setup ────────────────────────────────────────────────────────────────

Base.metadata.create_all(bind=engine)

app = FastAPI(title="Viva Question Module", version="1.0.0")
# Auth uses Authorization headers, not cookies, so wildcard origins are safe
# only with allow_credentials=False (browsers reject "*" + credentials anyway).
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_credentials=False,
                   allow_methods=["*"], allow_headers=["*"])


# ─── Auth Utilities ───────────────────────────────────────────────────────────

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="/auth/token")

def hash_password(password: str) -> str:
    return pwd_context.hash(password)

def verify_password(plain: str, hashed: str) -> bool:
    return pwd_context.verify(plain, hashed)

def create_token(teacher_id: int) -> str:
    expire = datetime.now(timezone.utc) + timedelta(minutes=settings.ACCESS_TOKEN_EXPIRE_MINUTES)
    return jwt.encode({"sub": str(teacher_id), "exp": expire}, settings.SECRET_KEY, algorithm=settings.ALGORITHM)

def get_current_teacher(token: str = Depends(oauth2_scheme), db: Session = Depends(get_db)):
    try:
        payload = jwt.decode(token, settings.SECRET_KEY, algorithms=[settings.ALGORITHM])
        sub_id = payload.get("sub")
        if sub_id is None:
            raise HTTPException(status_code=401, detail="Invalid token")
        teacher = db.query(Teacher).filter(Teacher.id == int(sub_id)).first()
        if not teacher:
            raise HTTPException(status_code=401, detail="Invalid token")
        return teacher
    except (JWTError, ValueError):
        raise HTTPException(status_code=401, detail="Invalid token")


# ─── File Parser ──────────────────────────────────────────────────────────────

def parse_file(filename: str, file_bytes: bytes) -> list:
    """Parse .pdf / .docx / .xlsx / .csv / .pptx → list of (question, answer) tuples."""
    ext = filename.rsplit(".", 1)[-1].lower()

    if ext == "pdf":
        import pdfplumber
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            text = "\n".join(p.extract_text() or "" for p in pdf.pages)
        return _parse_text(text)

    elif ext == "docx":
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        for table in doc.tables:
            result = _parse_table([[c.text for c in row.cells] for row in table.rows])
            if result: return result
        return _parse_text("\n".join(p.text for p in doc.paragraphs))

    elif ext == "xlsx":
        import openpyxl
        ws = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True).active
        return _parse_table([[str(c) if c else "" for c in row] for row in ws.iter_rows(values_only=True)])

    elif ext == "csv":
        import csv
        rows = list(csv.reader(io.StringIO(file_bytes.decode("utf-8", errors="replace"))))
        return _parse_table(rows)

    elif ext == "pptx":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        lines.append(text)
        return _parse_text("\n".join(lines))

    elif ext == "xls":
        raise HTTPException(status_code=400,
            detail="Legacy .xls files are not supported. Re-save the file as .xlsx.")

    else:
        raise HTTPException(status_code=400, detail=f"Unsupported file type .{ext}. Use pdf, docx, xlsx, pptx, or csv.")


def _parse_table(rows: list) -> list:
    """Find Question/Answer columns in a table."""
    if not rows: return []
    headers = [str(h).strip().lower() for h in rows[0]]
    try:
        q_idx = next(i for i, h in enumerate(headers) if "question" in h or h == "q")
        a_idx = next(i for i, h in enumerate(headers) if "answer" in h or h == "a")
    except StopIteration:
        raise HTTPException(status_code=400,
            detail="Columns named 'Question' and 'Answer' not found. Check your file headers.")
    pairs = [(re.sub(r"\s+", " ", str(r[q_idx]).strip()), re.sub(r"\s+", " ", str(r[a_idx]).strip()))
             for r in rows[1:] if len(r) > max(q_idx, a_idx)]
    return [(q, a) for q, a in pairs if q and a]


def _parse_text(text: str) -> list:
    """Parse Q:/A: formatted text."""
    matches = re.findall(
        r"^Q\d*[:.)]\s*(.+?)\s*^A\d*[:.)]\s*(.+?)(?=^Q\d*[:.)]|\Z)",
        text, re.DOTALL | re.IGNORECASE | re.MULTILINE
    )
    pairs = [
        (re.sub(r"\s+", " ", q.strip()), re.sub(r"\s+", " ", a.strip()))
        for q, a in matches if q.strip() and a.strip()
    ]
    if not pairs:
        raise HTTPException(status_code=400,
            detail="No Q&A found. Use 'Q: ...' / 'A: ...' format in Word/PDF files.")
    return pairs


def _validate_count(pairs: list) -> list:
    if len(pairs) < 10:
        raise HTTPException(status_code=400, detail=f"Only {len(pairs)} questions found. Minimum is 10.")
    if len(pairs) > 15:
        raise HTTPException(status_code=400, detail=f"{len(pairs)} questions found. Maximum is 15.")
    return pairs


# ─── Helper to build response dicts ──────────────────────────────────────────

def _set_response(s: QuestionSet, include_answers=False) -> dict:
    questions = sorted(s.questions, key=lambda q: q.question_number)
    q_list = [{"id": q.id, "question_number": q.question_number,
                "question_text": q.question_text,
                **({"answer_text": q.answer_text} if include_answers else {})}
              for q in questions]
    return {"id": s.id, "title": s.title, "subject": s.subject,
            "original_filename": s.original_filename, "created_at": s.created_at,
            "question_count": len(questions), "questions": q_list}

def _get_own_set(set_id: int, teacher_id: int, db: Session) -> QuestionSet:
    qs = db.query(QuestionSet).filter(QuestionSet.id == set_id,
                                      QuestionSet.teacher_id == teacher_id).first()
    if not qs: raise HTTPException(status_code=404, detail="Question set not found.")
    return qs


# ─── Auth Routes ──────────────────────────────────────────────────────────────

@app.post("/auth/register", response_model=TeacherOut, status_code=201, tags=["Auth"])
def register(payload: TeacherRegister, db: Session = Depends(get_db)):
    if db.query(Teacher).filter(Teacher.email == payload.email).first():
        raise HTTPException(status_code=400, detail="Email already registered.")
    teacher = Teacher(name=payload.name, email=payload.email,
                      hashed_password=hash_password(payload.password))
    db.add(teacher); db.commit(); db.refresh(teacher)
    return teacher

@app.post("/auth/login", response_model=Token, tags=["Auth"])
def login(payload: TeacherLogin, db: Session = Depends(get_db)):
    """Login using JSON body — use this for normal API calls."""
    teacher = db.query(Teacher).filter(Teacher.email == payload.email).first()
    if not teacher or not verify_password(payload.password, teacher.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password.")
    return {"access_token": create_token(teacher.id), "token_type": "bearer"}


@app.post("/auth/token", response_model=Token, include_in_schema=False)
def login_for_swagger(form_data: OAuth2PasswordRequestForm = Depends(), db: Session = Depends(get_db)):
    """OAuth2 form-based login — used internally by Swagger's Authorize button.
    Enter your email in the 'username' field and your password in 'password'."""
    teacher = db.query(Teacher).filter(Teacher.email == form_data.username).first()
    if not teacher or not verify_password(form_data.password, teacher.hashed_password):
        raise HTTPException(status_code=401, detail="Invalid email or password.")
    return {"access_token": create_token(teacher.id), "token_type": "bearer"}

@app.get("/auth/me", response_model=TeacherOut, tags=["Auth"])
def get_me(teacher: Teacher = Depends(get_current_teacher)):
    return teacher


# ─── Question Routes ──────────────────────────────────────────────────────────

@app.post("/questions/upload", response_model=QuestionSetDetailWithAnswers,
          status_code=201, tags=["Questions"])
async def upload_file(
    file: UploadFile = File(...),
    title: str = Form(...),
    subject: Optional[str] = Form(None),
    db: Session = Depends(get_db),
    teacher: Teacher = Depends(get_current_teacher),
):
    """Upload .pdf / .docx / .xlsx / .csv / .pptx with 10–15 Q&A pairs."""
    contents = await file.read()
    if len(contents) > 5 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="File too large. Max 5MB.")

    pairs = _validate_count(parse_file(file.filename, contents))

    qs = QuestionSet(title=title, subject=subject,
                     original_filename=file.filename, teacher_id=teacher.id)
    db.add(qs); db.flush()

    for idx, (q_text, a_text) in enumerate(pairs, start=1):
        db.add(Question(question_text=q_text, answer_text=a_text,
                        question_number=idx, question_set_id=qs.id))
    db.commit(); db.refresh(qs)
    return _set_response(qs, include_answers=True)


@app.get("/questions/my-sets", response_model=List[QuestionSetOut], tags=["Questions"])
def get_my_sets(db: Session = Depends(get_db), teacher: Teacher = Depends(get_current_teacher)):
    """List all question sets created by the logged-in teacher."""
    sets = db.query(QuestionSet).filter(QuestionSet.teacher_id == teacher.id)\
             .order_by(QuestionSet.created_at.desc()).all()
    return [{"id": s.id, "title": s.title, "subject": s.subject,
             "original_filename": s.original_filename, "created_at": s.created_at,
             "question_count": len(s.questions)} for s in sets]


@app.get("/questions/set/{set_id}/teacher", response_model=QuestionSetDetailWithAnswers, tags=["Questions"])
def get_set_teacher_view(set_id: int, db: Session = Depends(get_db),
                         teacher: Teacher = Depends(get_current_teacher)):
    """Teacher view — questions WITH answers."""
    return _set_response(_get_own_set(set_id, teacher.id, db), include_answers=True)


@app.get("/questions/set/{set_id}/viva", response_model=QuestionSetDetail, tags=["Questions"])
def get_set_viva_view(set_id: int, db: Session = Depends(get_db),
                      teacher: Teacher = Depends(get_current_teacher)):
    """Viva display — questions ONLY, no answers shown."""
    return _set_response(_get_own_set(set_id, teacher.id, db), include_answers=False)


@app.get("/questions/set/{set_id}/answers", tags=["Questions"])
def get_answers_for_evaluation(set_id: int, db: Session = Depends(get_db),
                                teacher: Teacher = Depends(get_current_teacher)):
    """For the Evaluation module — returns answer map keyed by question ID."""
    qs = _get_own_set(set_id, teacher.id, db)
    return {
        "question_set_id": qs.id,
        "answers": {q.id: {"question_number": q.question_number, "answer_text": q.answer_text}
                    for q in qs.questions}
    }


@app.delete("/questions/set/{set_id}", status_code=204, tags=["Questions"])
def delete_set(set_id: int, db: Session = Depends(get_db),
               teacher: Teacher = Depends(get_current_teacher)):
    """Delete a question set and all its questions."""
    qs = _get_own_set(set_id, teacher.id, db)
    db.delete(qs); db.commit()